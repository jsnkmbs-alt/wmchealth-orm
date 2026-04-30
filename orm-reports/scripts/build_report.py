"""
build_report.py
---------------
Reads metrics JSON from fetch_data.py and generates the
WMCHealth CTL PowerPoint report into orm-reports/docs/reports/.

Usage:
  python orm-reports/scripts/build_report.py --data orm-reports/scripts/data.json
"""

import os
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── COLORS ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)
TEAL   = RGBColor(0x0D, 0x7C, 0x7C)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xE6, 0x51, 0x00)
RED    = RGBColor(0xB7, 0x1C, 0x1C)
SLATE  = RGBColor(0x45, 0x5A, 0x64)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xF4, 0xF8, 0xFB)
LGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
LRED   = RGBColor(0xFF, 0xEB, 0xEE)
LBLUE  = RGBColor(0xE3, 0xF2, 0xFD)

REGION_COLORS = {"South": GREEN, "North": TEAL, "West": AMBER}


# ── HELPERS ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=11, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or SLATE
    return tb


def fmt_hours(h):
    if not h:
        return "N/A"
    if h < 1:
        return "<1 HR"
    if h < 24:
        return f"{int(round(h))} HRS"
    return f"{h/24:.1f} DAYS"


def eng_color(rate):
    if rate >= 90:
        return GREEN
    if rate >= 75:
        return AMBER
    return RED


def color_bar(slide, color, h=0.12):
    add_rect(slide, 0, 0, 10, h, color)


# ── SLIDE 1 — TITLE ───────────────────────────────────────────────────────────
def slide_title(prs, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_rect(s, 0.6, 3.1, 2.5, 0.05, RGBColor(0x7E, 0xC8, 0xE3))
    add_text(s, "CLOSING THE LOOP", 0.6, 0.75, 9, 0.4,
             size=11, bold=True, color=RGBColor(0x7E, 0xC8, 0xE3))
    add_text(s, "Engagement Report", 0.6, 1.2, 8.5, 1.1,
             size=40, bold=True, color=WHITE)
    add_text(s, f"{meta['month_label']}  |  Westchester Medical Center Health Network",
             0.6, 2.5, 9, 0.5, size=14, color=RGBColor(0xA0, 0xC4, 0xD8))
    add_text(s, f"Source: Press Ganey  |  {meta['from_date'][:10]} – {meta['to_date'][:10]}  |  Generated: {meta['generated_at'][:10]}",
             0.6, 4.8, 9, 0.35, size=9, color=RGBColor(0x6A, 0x9B, 0xB5))


# ── SLIDE 2 — NETWORK OVERVIEW ────────────────────────────────────────────────
def slide_overview(prs, metrics, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFFWHT)
    net = metrics["network"]

    add_text(s, f"Network Overview – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
             size=24, bold=True, color=NAVY)
    add_text(s, "All Regions  |  Westchester Medical Center Health Network",
             0.4, 0.75, 9.2, 0.3, size=10, color=SLATE)

    cards = [
        ("Avg Engagement Rate", f"{net['engagement_rate']}%",  "vs prior year",          TEAL),
        ("Total Tasks",          str(net["total_tasks"]),       "total feedback records",  NAVY),
        ("Closed Tasks",         str(net["closed_tasks"]),      f"{net['engagement_rate']}% close rate", GREEN),
        ("Avg Time to Close",    fmt_hours(net["avg_time_to_close_hours"]), "from review to response", AMBER),
    ]
    for i, (label, val, sub, color) in enumerate(cards):
        x = 0.3 + i * 2.35
        add_rect(s, x, 1.2, 2.15, 1.5, WHITE, LGRAY)
        add_rect(s, x, 1.2, 2.15, 0.08, color)
        add_text(s, val,   x+0.05, 1.32, 2.05, 0.65, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, label, x+0.05, 1.95, 2.05, 0.3,  size=9,  bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(s, sub,   x+0.05, 2.25, 2.05, 0.3,  size=8,  color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER)

    # Sentiment bar
    total = net["positive"] + net["neutral"] + net["negative"]
    if total:
        pp = round(net["positive"] / total * 100)
        np_ = round(net["neutral"]  / total * 100)
        ng  = 100 - pp - np_
        add_text(s, "Sentiment Distribution", 0.4, 3.0, 5, 0.3, size=11, bold=True, color=NAVY)
        bw, bx, by, bh = 9.2, 0.4, 3.35, 0.45
        add_rect(s, bx,               by, bw*pp/100,  bh, GREEN)
        add_rect(s, bx+bw*pp/100,     by, bw*np_/100, bh, AMBER)
        add_rect(s, bx+bw*(pp+np_)/100, by, bw*ng/100, bh, RED)
        add_text(s, f"✔ Positive {pp}%",  bx,       by+0.5, 3, 0.3, size=10, bold=True, color=GREEN)
        add_text(s, f"◆ Neutral {np_}%",  bx+3.1,   by+0.5, 3, 0.3, size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_text(s, f"✖ Negative {ng}%",  bx+6.2,   by+0.5, 3, 0.3, size=10, bold=True, color=RED,   align=PP_ALIGN.RIGHT)

    # Top locations
    add_text(s, "Top Locations by Engagement", 0.4, 4.15, 9.2, 0.3, size=11, bold=True, color=NAVY)
    locs = sorted([(l, d) for l, d in metrics["by_location"].items() if d["total"] > 0],
                  key=lambda x: x[1]["engagement_rate"], reverse=True)[:6]
    for i, (loc, d) in enumerate(locs):
        ry = 4.5 + i * 0.32
        bg = OFFWHT if i % 2 == 0 else WHITE
        add_rect(s, 0.4, ry, 9.2, 0.3, bg, LGRAY)
        short = (loc[:45] + "...") if len(loc) > 45 else loc
        add_text(s, short, 0.5, ry+0.04, 6, 0.24, size=9, color=SLATE)
        add_text(s, f"{d['engagement_rate']}%", 6.6, ry+0.04, 1.2, 0.24,
                 size=9, bold=True, color=eng_color(d["engagement_rate"]), align=PP_ALIGN.RIGHT)
        add_text(s, d["region"], 8.0, ry+0.04, 1.5, 0.24,
                 size=9, color=REGION_COLORS.get(d["region"], SLATE), align=PP_ALIGN.RIGHT)


# ── SLIDE 3 — REGIONAL OVERVIEW ───────────────────────────────────────────────
def slide_regional(prs, metrics, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFFWHT)
    add_text(s, f"Regional Performance – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
             size=24, bold=True, color=NAVY)
    add_text(s, "Engagement Rate, Tasks & Avg Time to Close by Region",
             0.4, 0.75, 9.2, 0.3, size=10, color=SLATE)

    for i, region in enumerate(["South", "North", "West"]):
        rm    = metrics["regions"].get(region, {})
        color = REGION_COLORS[region]
        x     = 0.35 + i * 3.1
        rate  = rm.get("engagement_rate", 0)
        add_rect(s, x, 1.2, 2.95, 2.1, WHITE, LGRAY)
        add_rect(s, x, 1.2, 2.95, 0.45, color)
        add_text(s, region, x+0.1, 1.22, 2.2, 0.4, size=14, bold=True, color=WHITE)
        add_text(s, f"{rate}%", x+0.05, 1.68, 2.85, 0.65,
                 size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, f"{rm.get('total_tasks',0)} Tasks  |  Avg: {fmt_hours(rm.get('avg_time_to_close_hours',0))}",
                 x+0.05, 2.35, 2.85, 0.3, size=9, color=SLATE, align=PP_ALIGN.CENTER)
        total = rm.get("positive",0) + rm.get("neutral",0) + rm.get("negative",0)
        if total:
            pw = 2.85 * rm["positive"] / total
            nw = 2.85 * rm["neutral"]  / total
            gw = 2.85 - pw - nw
            add_rect(s, x+0.05, 2.75, pw, 0.2, GREEN)
            add_rect(s, x+0.05+pw, 2.75, nw, 0.2, AMBER)
            add_rect(s, x+0.05+pw+nw, 2.75, gw, 0.2, RED)

    # Source table
    add_text(s, "Engagement by Source", 0.4, 3.55, 9.2, 0.3, size=11, bold=True, color=NAVY)
    cols  = ["Source", "Total", "Closed", "Rate"]
    cxs   = [0.4, 4.5, 6.0, 7.5]
    cws   = [4.0, 1.4, 1.4, 1.6]
    add_rect(s, 0.4, 3.9, 9.2, 0.32, NAVY)
    for j, h in enumerate(cols):
        add_text(s, h, cxs[j]+0.05, 3.92, cws[j], 0.28, size=9, bold=True, color=WHITE)
    for i, (src, sd) in enumerate(sorted(metrics["by_source"].items(), key=lambda x: x[1]["total"], reverse=True)):
        ry = 4.23 + i * 0.32
        add_rect(s, 0.4, ry, 9.2, 0.3, OFFWHT if i%2==0 else WHITE, LGRAY)
        add_text(s, src,               cxs[0]+0.05, ry+0.04, cws[0], 0.24, size=9, color=SLATE)
        add_text(s, str(sd["total"]),  cxs[1]+0.05, ry+0.04, cws[1], 0.24, size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(s, str(sd["closed"]), cxs[2]+0.05, ry+0.04, cws[2], 0.24, size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(s, f"{sd['engagement_rate']}%", cxs[3]+0.05, ry+0.04, cws[3], 0.24,
                 size=9, bold=True, color=eng_color(sd["engagement_rate"]), align=PP_ALIGN.CENTER)


# ── SLIDE 4 — TASK OWNER ─────────────────────────────────────────────────────
def slide_owners(prs, metrics, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFFWHT)
    add_text(s, f"Task Owner Performance – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
             size=24, bold=True, color=NAVY)
    add_text(s, "Engagement rates by task owner across all regions",
             0.4, 0.75, 9.2, 0.3, size=10, color=SLATE)

    cols = ["Owner", "Total", "Closed", "Engagement Rate", "Avg Time to Close"]
    cxs  = [0.4, 3.8, 5.2, 6.5, 8.2]
    cws  = [3.3, 1.3, 1.2, 1.6, 1.5]
    add_rect(s, 0.4, 1.15, 9.2, 0.35, NAVY)
    for j, h in enumerate(cols):
        add_text(s, h, cxs[j]+0.05, 1.18, cws[j], 0.3, size=9, bold=True, color=WHITE)

    owners = sorted(metrics["by_owner"].items(), key=lambda x: x[1]["engagement_rate"], reverse=True)
    for i, (owner, od) in enumerate(owners):
        ry = 1.52 + i * 0.35
        add_rect(s, 0.4, ry, 9.2, 0.33, OFFWHT if i%2==0 else WHITE, LGRAY)
        short = (owner[:40] + "...") if len(owner) > 40 else owner
        add_text(s, short,                      cxs[0]+0.05, ry+0.05, cws[0], 0.25, size=9, color=SLATE)
        add_text(s, str(od["total"]),            cxs[1]+0.05, ry+0.05, cws[1], 0.25, size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(s, str(od["closed"]),           cxs[2]+0.05, ry+0.05, cws[2], 0.25, size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(s, f"{od['engagement_rate']}%", cxs[3]+0.05, ry+0.05, cws[3], 0.25,
                 size=9, bold=True, color=eng_color(od["engagement_rate"]), align=PP_ALIGN.CENTER)
        add_text(s, fmt_hours(od["avg_time_to_close_hours"]), cxs[4]+0.05, ry+0.05, cws[4], 0.25,
                 size=9, color=SLATE, align=PP_ALIGN.CENTER)

    zero = [o for o, d in metrics["by_owner"].items() if d["engagement_rate"] == 0 and d["total"] > 0]
    if zero:
        wy = 1.52 + len(owners) * 0.35 + 0.15
        add_rect(s, 0.4, wy, 9.2, 0.4, LRED, RED)
        add_text(s, f"⚠  {', '.join(zero)} {'has' if len(zero)==1 else 'have'} 0% engagement — "
                    f"{sum(metrics['by_owner'][o]['total'] for o in zero)} tasks unresolved network-wide.",
                 0.55, wy+0.07, 9.0, 0.28, size=9, bold=True, color=RED)


# ── SLIDES 5 & 6-8 — ALERTS (network + per region) ───────────────────────────
def slide_alerts(prs, metrics, meta, region=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFFWHT)

    if region:
        color_bar(s, REGION_COLORS.get(region, NAVY))
        title = f"Critical Alerts – {region} Region"
    else:
        title = "Critical Alerts & Escalations"

    add_text(s, title, 0.4, 0.25, 9.2, 0.5, size=24, bold=True, color=NAVY)
    add_text(s, f"High-priority items requiring leadership attention  |  {meta['month_label']}",
             0.4, 0.75, 9.2, 0.28, size=10, color=SLATE)

    alerts = metrics["critical_alerts"]
    if region:
        alerts = [a for a in alerts if a.get("region") == region]

    if not alerts:
        add_rect(s, 0.4, 1.3, 9.2, 1.0, LGREEN, RGBColor(0xA5, 0xD6, 0xA7))
        add_text(s, f"✔  No critical alerts for {'this region' if region else 'the network'} this month.",
                 0.6, 1.65, 9.0, 0.4, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        return

    top = sorted(alerts, key=lambda x: (x.get("score", 5), -x.get("days_open", 0)))[:5]
    for i, a in enumerate(top):
        y     = 1.1 + i * 0.9
        score = a.get("score", 0)
        days  = a.get("days_open", 0)
        bg    = LRED if score < 2 else RGBColor(0xFF, 0xF3, 0xE0)
        clr   = RED  if score < 2 else AMBER
        add_rect(s, 0.4, y, 9.2, 0.82, bg, clr)
        add_rect(s, 0.4, y, 0.08, 0.82, clr)
        title_t = f"{a.get('location','?')}  |  PFS: {score:.1f}  |  {days} days open  |  Owner: {a.get('owner','?')}  |  {a.get('source','')}"
        excerpt = (a.get("text") or "")[:120] + ("..." if len(a.get("text") or "") > 120 else "")
        add_text(s, title_t, 0.55, y+0.05, 8.9, 0.26, size=9.5, bold=True, color=clr)
        add_text(s, excerpt or "(No review text)", 0.55, y+0.33, 8.9, 0.44, size=8.5, color=SLATE)

    if len(alerts) > 5:
        add_text(s, f"+ {len(alerts)-5} additional open tasks — see full data export.",
                 0.4, 1.1+5*0.9+0.05, 9.2, 0.28, size=9, color=SLATE, italic=True)


# ── SLIDES 9 & 10-12 — RECOMMENDATIONS (network + per region) ────────────────
def slide_recs(prs, metrics, meta, region=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    if region:
        color_bar(s, REGION_COLORS.get(region, TEAL))
        title = f"Recommendations – {region} Region"
    else:
        title = "Recommendations & Next Steps"

    add_text(s, title, 0.5, 0.25, 9, 0.55, size=24, bold=True, color=WHITE)
    add_rect(s, 0.5, 0.85, 1.5, 0.05, RGBColor(0x7E, 0xC8, 0xE3))

    alerts = metrics["critical_alerts"]
    if region:
        alerts = [a for a in alerts if a.get("region") == region]

    zero_owners  = [o for o, d in metrics["by_owner"].items() if d["engagement_rate"] == 0 and d["total"] > 0]
    aged         = [a for a in alerts if a.get("days_open", 0) >= 14]
    critical_scr = [a for a in alerts if a.get("score", 5) < 1.5]
    low_sources  = [(s_, d) for s_, d in metrics["by_source"].items() if d["engagement_rate"] < 50 and d["total"] > 0]

    recs = []
    rec_colors = [RGBColor(0xEF,0x53,0x50), RGBColor(0xFF,0xA7,0x26),
                  RGBColor(0x42,0xA5,0xF5), RGBColor(0x66,0xBB,0x6A), RGBColor(0xAB,0x47,0xBC)]

    if aged:
        locs = list(set(a["location"] for a in aged[:3]))
        recs.append({"title": f"Close {len(aged)} Aged Open Tasks (14+ days)",
                     "body":  f"Prioritize: {', '.join(locs)}. Target closure within 24 hours."})
    if critical_scr:
        recs.append({"title": f"Escalate {len(critical_scr)} Critical Score Reviews (PFS < 1.5)",
                     "body":  "Requires direct patient outreach, clinical review, and Patient Experience follow-up."})
    if zero_owners:
        recs.append({"title": "Recover 0% Engagement Owner Groups",
                     "body":  f"{', '.join(zero_owners[:2])} at 0% engagement. Clarify ownership and set daily SLA review."})
    if low_sources:
        names = ", ".join(n for n, _ in low_sources)
        recs.append({"title": f"Improve Response Rate on {names}",
                     "body":  "Ensure all platforms are included in the managed engagement workflow."})
    recs.append({"title": "Celebrate Positive Staff Feedback",
                 "body":  "Share named staff mentions with department heads. Nominate for Care Champions program."})

    dark_bg  = RGBColor(0x0A, 0x20, 0x40)
    dark_bdr = RGBColor(0x1A, 0x3F, 0x6F)
    cols = 2 if len(recs) > 2 else 1
    w    = 4.5 if cols == 2 else 9.0

    for i, r in enumerate(recs[:6]):
        col = i % cols
        row = i // cols
        x = 0.5 + col * (w + 0.2)
        y = 1.05 + row * 1.5
        clr = rec_colors[i % len(rec_colors)]
        add_rect(s, x, y, w, 1.35, dark_bg, dark_bdr)
        add_text(s, f"0{i+1}", x+0.12, y+0.1, 0.55, 0.55, size=20, bold=True, color=clr)
        add_text(s, r["title"], x+0.72, y+0.1,  w-0.88, 0.38, size=10,  bold=True, color=WHITE)
        add_text(s, r["body"],  x+0.72, y+0.5,  w-0.88, 0.76, size=8.5, color=RGBColor(0xA0,0xC4,0xD8))


# ── MAIN BUILD ────────────────────────────────────────────────────────────────
def build(metrics: dict, output_path: str):
    meta = metrics["meta"]
    prs  = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)

    print("🖼️  Building slides...")
    slide_title(prs, meta)
    slide_overview(prs, metrics, meta)
    slide_regional(prs, metrics, meta)
    slide_owners(prs, metrics, meta)

    slide_alerts(prs, metrics, meta)
    for region in ["South", "North", "West"]:
        slide_alerts(prs, metrics, meta, region=region)

    slide_recs(prs, metrics, meta)
    for region in ["South", "North", "West"]:
        slide_recs(prs, metrics, meta, region=region)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ Report saved → {output_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data",   default="orm-reports/scripts/data.json")
    parser.add_argument("--output", default=None)
    args = parser.parse_args()

    with open(args.data) as f:
        metrics = json.load(f)

    month_str   = metrics["meta"]["from_date"][:7]
    output_path = args.output or f"orm-reports/docs/reports/WMCHealth_CTL_{month_str}.pptx"
    build(metrics, output_path)


if __name__ == "__main__":
    main()
