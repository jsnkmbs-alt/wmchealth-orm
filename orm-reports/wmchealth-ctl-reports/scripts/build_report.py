"""
build_report.py
---------------
Reads the metrics JSON produced by fetch_data.py and generates
the WMCHealth CTL PowerPoint report, then saves it to docs/reports/.

Usage:
  python scripts/build_report.py --data scripts/data.json
"""

import os
import json
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# ── COLORS ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)
TEAL   = RGBColor(0x0D, 0x7C, 0x7C)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xE6, 0x51, 0x00)
RED    = RGBColor(0xB7, 0x1C, 0x1C)
SLATE  = RGBColor(0x45, 0x5A, 0x64)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xF4, 0xF8, 0xFB)
LGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
LRED   = RGBColor(0xFF, 0xEB, 0xEE)
LBLUE  = RGBColor(0xE3, 0xF2, 0xFD)
REGION_COLORS = {
    "South": GREEN,
    "North": TEAL,
    "West":  AMBER,
}


# ── SLIDE HELPERS ─────────────────────────────────────────────────────────────
def inches(*args):
    return [Inches(a) for a in args]


def add_textbox(slide, text, x, y, w, h, font_size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, bg_color=None,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or SLATE
    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = PtU(line_width)
    else:
        shape.line.fill.background()
    return shape


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_colored_bar(slide, color, height=0.12):
    """Add a thin colored bar at the very top of a slide."""
    add_rect(slide, 0, 0, 10, height, color)


def engagement_color(rate):
    if rate >= 90:
        return GREEN
    elif rate >= 75:
        return AMBER
    else:
        return RED


def fmt_hours(hours):
    if hours < 1:
        return "<1 HR"
    elif hours < 24:
        return f"{int(round(hours))} HRS"
    else:
        days = hours / 24
        return f"{days:.1f} DAYS"


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def build_title_slide(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    add_textbox(slide, "CLOSING THE LOOP", 0.6, 0.75, 9, 0.4,
                font_size=11, bold=True, color=RGBColor(0x7E, 0xC8, 0xE3),
                align=PP_ALIGN.LEFT)
    add_textbox(slide, "Engagement Report", 0.6, 1.2, 8.5, 1.2,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, f"{meta['month_label']}  |  Westchester Medical Center Health Network",
                0.6, 2.5, 9, 0.5, font_size=14, color=RGBColor(0xA0, 0xC4, 0xD8),
                align=PP_ALIGN.LEFT)
    add_rect(slide, 0.6, 3.1, 2.5, 0.05, RGBColor(0x7E, 0xC8, 0xE3))
    add_textbox(slide, f"Source: Press Ganey  |  Report Period: {meta['from_date'][:10]} – {meta['to_date'][:10]}  |  Generated: {meta['generated_at'][:10]}",
                0.6, 4.8, 9, 0.35, font_size=9,
                color=RGBColor(0x6A, 0x9B, 0xB5), align=PP_ALIGN.LEFT)


def build_overview_slide(prs, metrics, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, OFFWHT)
    net = metrics["network"]

    add_textbox(slide, f"Network Overview – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
                font_size=24, bold=True, color=NAVY)
    add_textbox(slide, "All Regions  |  Westchester Medical Center Health Network",
                0.4, 0.75, 9.2, 0.3, font_size=10, color=SLATE)

    # KPI cards
    cards = [
        ("Avg Engagement Rate", f"{net['engagement_rate']}%", "vs prior year", TEAL),
        ("Total Tasks", str(net["total_tasks"]), "total feedback records", NAVY),
        ("Online/Email Closed", str(net["closed_tasks"]), f"{net['engagement_rate']}% close rate", GREEN),
        ("Avg Time to Close", fmt_hours(net["avg_time_to_close_hours"]), "same business day", AMBER),
    ]
    for i, (label, val, sub, color) in enumerate(cards):
        x = 0.3 + i * 2.35
        add_rect(slide, x, 1.2, 2.15, 1.5, WHITE, LGRAY, 0.5)
        add_rect(slide, x, 1.2, 2.15, 0.08, color)
        add_textbox(slide, val, x + 0.05, 1.32, 2.05, 0.65,
                    font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x + 0.05, 1.95, 2.05, 0.3,
                    font_size=9, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, x + 0.05, 2.25, 2.05, 0.3,
                    font_size=8, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER)

    # Sentiment summary bar
    total = net["positive"] + net["neutral"] + net["negative"]
    if total:
        pos_pct = round(net["positive"] / total * 100)
        neu_pct = round(net["neutral"] / total * 100)
        neg_pct = 100 - pos_pct - neu_pct

        add_textbox(slide, "Sentiment Distribution (All Reviews)", 0.4, 3.0, 5, 0.3,
                    font_size=11, bold=True, color=NAVY)
        bar_w = 9.2
        bar_x = 0.4
        bar_y = 3.35
        bar_h = 0.45
        add_rect(slide, bar_x, bar_y, bar_w * pos_pct / 100, bar_h, GREEN)
        add_rect(slide, bar_x + bar_w * pos_pct / 100, bar_y, bar_w * neu_pct / 100, bar_h, AMBER)
        add_rect(slide, bar_x + bar_w * (pos_pct + neu_pct) / 100, bar_y, bar_w * neg_pct / 100, bar_h, RED)
        add_textbox(slide, f"✔ Positive {pos_pct}%", bar_x, bar_y + 0.5, 3, 0.3,
                    font_size=10, bold=True, color=GREEN)
        add_textbox(slide, f"◆ Neutral {neu_pct}%", bar_x + 3.1, bar_y + 0.5, 3, 0.3,
                    font_size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"✖ Negative {neg_pct}%", bar_x + 6.2, bar_y + 0.5, 3, 0.3,
                    font_size=10, bold=True, color=RED, align=PP_ALIGN.RIGHT)

    # Top locations table
    add_textbox(slide, "Top Locations by Engagement", 0.4, 4.15, 9.2, 0.3,
                font_size=11, bold=True, color=NAVY)

    sorted_locs = sorted(
        [(loc, d) for loc, d in metrics["by_location"].items() if d["total"] > 0],
        key=lambda x: x[1]["engagement_rate"], reverse=True
    )[:6]

    for i, (loc, d) in enumerate(sorted_locs):
        row_y = 4.5 + i * 0.32
        bg = RGBColor(0xF4, 0xF8, 0xFB) if i % 2 == 0 else WHITE
        add_rect(slide, 0.4, row_y, 9.2, 0.3, bg, LGRAY, 0.3)
        short_name = loc[:45] + "..." if len(loc) > 45 else loc
        add_textbox(slide, short_name, 0.5, row_y + 0.04, 6, 0.24, font_size=9, color=SLATE)
        rate = d["engagement_rate"]
        add_textbox(slide, f"{rate}%", 6.6, row_y + 0.04, 1.2, 0.24,
                    font_size=9, bold=True, color=engagement_color(rate), align=PP_ALIGN.RIGHT)
        add_textbox(slide, d["region"], 8.0, row_y + 0.04, 1.5, 0.24,
                    font_size=9, color=REGION_COLORS.get(d["region"], SLATE), align=PP_ALIGN.RIGHT)


def build_regional_overview_slide(prs, metrics, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, OFFWHT)

    add_textbox(slide, f"Regional Performance – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
                font_size=24, bold=True, color=NAVY)
    add_textbox(slide, "Engagement Rate, Total Tasks & Avg Time to Close by Region",
                0.4, 0.75, 9.2, 0.3, font_size=10, color=SLATE)

    regions = [
        ("South", "▲ Best Rate"),
        ("North", "Fastest Close"),
        ("West", ""),
    ]
    for i, (region, badge) in enumerate(regions):
        rm = metrics["regions"].get(region, {})
        color = REGION_COLORS.get(region, SLATE)
        x = 0.35 + i * 3.1
        add_rect(slide, x, 1.2, 2.95, 2.1, WHITE, LGRAY, 0.5)
        add_rect(slide, x, 1.2, 2.95, 0.45, color)
        add_textbox(slide, region, x + 0.1, 1.22, 2.0, 0.4,
                    font_size=14, bold=True, color=WHITE)
        if badge:
            add_textbox(slide, badge, x + 1.8, 1.28, 1.1, 0.28,
                        font_size=7, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        rate = rm.get("engagement_rate", 0)
        add_textbox(slide, f"{rate}%", x + 0.05, 1.68, 2.85, 0.65,
                    font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{rm.get('total_tasks', 0)} Tasks  |  Avg Close: {fmt_hours(rm.get('avg_time_to_close_hours', 0))}",
                    x + 0.05, 2.35, 2.85, 0.3,
                    font_size=9, color=SLATE, align=PP_ALIGN.CENTER)

        # Mini sentiment bar
        total = rm.get("positive", 0) + rm.get("neutral", 0) + rm.get("negative", 0)
        if total:
            p_w = 2.85 * rm.get("positive", 0) / total
            n_w = 2.85 * rm.get("neutral", 0) / total
            neg_w = 2.85 - p_w - n_w
            bar_y = 2.75
            add_rect(slide, x + 0.05, bar_y, p_w, 0.2, GREEN)
            add_rect(slide, x + 0.05 + p_w, bar_y, n_w, 0.2, AMBER)
            add_rect(slide, x + 0.05 + p_w + n_w, bar_y, neg_w, 0.2, RED)

    # Source breakdown table
    add_textbox(slide, "Engagement by Source", 0.4, 3.5, 9.2, 0.3,
                font_size=11, bold=True, color=NAVY)

    headers = ["Source", "Total Reviews", "Engaged", "Response Rate"]
    col_x = [0.4, 3.5, 5.5, 7.5]
    col_w = [3.0, 1.8, 1.8, 1.8]
    add_rect(slide, 0.4, 3.85, 9.2, 0.32, NAVY)
    for j, h in enumerate(headers):
        add_textbox(slide, h, col_x[j] + 0.05, 3.87, col_w[j], 0.28,
                    font_size=9, bold=True, color=WHITE)

    for i, (src, sd) in enumerate(sorted(metrics["by_source"].items(), key=lambda x: x[1]["total"], reverse=True)):
        row_y = 4.18 + i * 0.32
        bg = OFFWHT if i % 2 == 0 else WHITE
        add_rect(slide, 0.4, row_y, 9.2, 0.3, bg, LGRAY, 0.3)
        add_textbox(slide, src, col_x[0] + 0.05, row_y + 0.04, col_w[0], 0.24, font_size=9, color=SLATE)
        add_textbox(slide, str(sd["total"]), col_x[1] + 0.05, row_y + 0.04, col_w[1], 0.24, font_size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_textbox(slide, str(sd["closed"]), col_x[2] + 0.05, row_y + 0.04, col_w[2], 0.24, font_size=9, color=SLATE, align=PP_ALIGN.CENTER)
        rate = sd["engagement_rate"]
        add_textbox(slide, f"{rate}%", col_x[3] + 0.05, row_y + 0.04, col_w[3], 0.24,
                    font_size=9, bold=True, color=engagement_color(rate), align=PP_ALIGN.CENTER)


def build_owner_slide(prs, metrics, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, OFFWHT)

    add_textbox(slide, f"Task Owner Performance – {meta['month_label']}", 0.4, 0.2, 9.2, 0.55,
                font_size=24, bold=True, color=NAVY)
    add_textbox(slide, "Engagement rates by task owner across all regions",
                0.4, 0.75, 9.2, 0.3, font_size=10, color=SLATE)

    headers = ["Owner", "Total Tasks", "Closed", "Engagement Rate", "Avg Time to Close"]
    col_x = [0.4, 3.8, 5.2, 6.5, 8.2]
    col_w = [3.3, 1.3, 1.2, 1.6, 1.5]

    add_rect(slide, 0.4, 1.15, 9.2, 0.35, NAVY)
    for j, h in enumerate(headers):
        add_textbox(slide, h, col_x[j] + 0.05, 1.18, col_w[j], 0.3,
                    font_size=9, bold=True, color=WHITE)

    sorted_owners = sorted(metrics["by_owner"].items(), key=lambda x: x[1]["engagement_rate"], reverse=True)

    for i, (owner, od) in enumerate(sorted_owners):
        row_y = 1.52 + i * 0.35
        bg = OFFWHT if i % 2 == 0 else WHITE
        add_rect(slide, 0.4, row_y, 9.2, 0.33, bg, LGRAY, 0.3)
        short = owner[:40] + "..." if len(owner) > 40 else owner
        add_textbox(slide, short, col_x[0] + 0.05, row_y + 0.05, col_w[0], 0.25, font_size=9, color=SLATE)
        add_textbox(slide, str(od["total"]), col_x[1] + 0.05, row_y + 0.05, col_w[1], 0.25,
                    font_size=9, color=SLATE, align=PP_ALIGN.CENTER)
        add_textbox(slide, str(od["closed"]), col_x[2] + 0.05, row_y + 0.05, col_w[2], 0.25,
                    font_size=9, color=SLATE, align=PP_ALIGN.CENTER)
        rate = od["engagement_rate"]
        add_textbox(slide, f"{rate}%", col_x[3] + 0.05, row_y + 0.05, col_w[3], 0.25,
                    font_size=9, bold=True, color=engagement_color(rate), align=PP_ALIGN.CENTER)
        add_textbox(slide, fmt_hours(od["avg_time_to_close_hours"]), col_x[4] + 0.05, row_y + 0.05, col_w[4], 0.25,
                    font_size=9, color=SLATE, align=PP_ALIGN.CENTER)

    # Warning if any owner at 0%
    zero_owners = [o for o, d in metrics["by_owner"].items() if d["engagement_rate"] == 0 and d["total"] > 0]
    if zero_owners:
        warn_y = 1.52 + len(sorted_owners) * 0.35 + 0.2
        add_rect(slide, 0.4, warn_y, 9.2, 0.4, LRED, RED, 0.5)
        add_textbox(slide, f"⚠  {', '.join(zero_owners)} {'has' if len(zero_owners)==1 else 'have'} 0% engagement — {sum(metrics['by_owner'][o]['total'] for o in zero_owners)} open tasks remain unresolved network-wide.",
                    0.55, warn_y + 0.06, 9.0, 0.3, font_size=9, bold=True, color=RED)


def build_alerts_slide(prs, metrics, meta, region=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, OFFWHT)

    if region:
        color = REGION_COLORS.get(region, NAVY)
        add_colored_bar(slide, color)
        title = f"Critical Alerts – {region} Region"
    else:
        title = "Critical Alerts & Escalations"

    add_textbox(slide, title, 0.4, 0.25, 9.2, 0.5, font_size=24, bold=True, color=NAVY)
    add_textbox(slide, f"High-priority items requiring leadership attention  |  {meta['month_label']}",
                0.4, 0.75, 9.2, 0.28, font_size=10, color=SLATE)

    # Filter alerts by region if needed
    alerts = metrics["critical_alerts"]
    if region:
        alerts = [a for a in alerts if a.get("region") == region]

    if not alerts:
        add_rect(slide, 0.4, 1.3, 9.2, 1.0, LGREEN, RGBColor(0xA5, 0xD6, 0xA7), 0.5)
        add_textbox(slide, f"✔  No critical alerts for {'this region' if region else 'the network'} this month.",
                    0.6, 1.6, 9.0, 0.4, font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        return

    # Sort by score ascending (worst first), then days open descending
    alerts_sorted = sorted(alerts, key=lambda x: (x.get("score", 5), -x.get("days_open", 0)))

    for i, alert in enumerate(alerts_sorted[:5]):
        y = 1.1 + i * 0.9
        score = alert.get("score", 0)
        days = alert.get("days_open", 0)
        bg = LRED if score < 2 else RGBColor(0xFF, 0xF3, 0xE0)
        border = RED if score < 2 else AMBER
        bar_color = RED if score < 2 else AMBER

        add_rect(slide, 0.4, y, 9.2, 0.82, bg, border, 0.5)
        add_rect(slide, 0.4, y, 0.08, 0.82, bar_color)

        location = alert.get("location", "Unknown location")
        owner = alert.get("owner", "Unassigned")
        source = alert.get("source", "")
        text = (alert.get("text") or "")[:120] + ("..." if len(alert.get("text") or "") > 120 else "")

        title_text = f"{location}  |  PFS: {score:.1f}  |  {days} days open  |  Owner: {owner}  |  {source}"
        add_textbox(slide, title_text, 0.55, y + 0.05, 8.9, 0.26,
                    font_size=9.5, bold=True, color=bar_color)
        add_textbox(slide, text or "(No review text provided)", 0.55, y + 0.32, 8.9, 0.44,
                    font_size=8.5, color=SLATE)

    if len(alerts) > 5:
        add_textbox(slide, f"+ {len(alerts) - 5} additional open tasks not shown — see full data export.",
                    0.4, 1.1 + 5 * 0.9 + 0.05, 9.2, 0.28, font_size=9, color=SLATE, italic=True)


def build_recommendations_slide(prs, metrics, meta, region=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    if region:
        color = REGION_COLORS.get(region, TEAL)
        add_colored_bar(slide, color)
        title = f"Recommendations – {region} Region"
    else:
        title = "Recommendations & Next Steps"

    add_textbox(slide, title, 0.5, 0.25, 9, 0.55, font_size=24, bold=True, color=WHITE)
    add_rect(slide, 0.5, 0.85, 1.5, 0.05, RGBColor(0x7E, 0xC8, 0xE3))

    # Auto-generate recommendations from the data
    recs = []
    alerts = metrics["critical_alerts"]
    if region:
        alerts = [a for a in alerts if a.get("region") == region]

    zero_owners = [o for o, d in metrics["by_owner"].items() if d["engagement_rate"] == 0 and d["total"] > 0]
    aged_tasks = [a for a in alerts if a.get("days_open", 0) >= 14]
    critical_score = [a for a in alerts if a.get("score", 5) < 1.5]
    low_sources = [(s, d) for s, d in metrics["by_source"].items() if d["engagement_rate"] < 50 and d["total"] > 0]

    if aged_tasks:
        locs = list(set(a["location"] for a in aged_tasks[:3]))
        recs.append({
            "num": "01", "color": RGBColor(0xEF, 0x53, 0x50),
            "title": f"Close {len(aged_tasks)} Aged Open Tasks (14+ days)",
            "body": f"Assign immediate attention to open tasks at: {', '.join(locs)}. Target closure within 24 hours."
        })
    if critical_score:
        recs.append({
            "num": "02", "color": RGBColor(0xFF, 0xA7, 0x26),
            "title": f"Escalate {len(critical_score)} Critical Score Reviews (PFS < 1.5)",
            "body": f"Reviews with the lowest PFS scores require direct patient outreach and clinical review. Contact Patient Experience."
        })
    if zero_owners:
        recs.append({
            "num": "03", "color": RGBColor(0x42, 0xA5, 0xF5),
            "title": "Recover 0% Engagement Owner Groups",
            "body": f"{', '.join(zero_owners[:2])} {'have' if len(zero_owners) > 1 else 'has'} 0% engagement. Clarify ownership, set SLA, implement daily review cadence."
        })
    if low_sources:
        src_names = ", ".join(s for s, _ in low_sources)
        recs.append({
            "num": "04", "color": RGBColor(0x66, 0xBB, 0x6A),
            "title": f"Improve Response Rate on {src_names}",
            "body": f"Engagement rate on {src_names} is below 50%. Ensure all reviews on every platform are included in the managed engagement workflow."
        })

    # Always add best-practice recommendation
    recs.append({
        "num": f"0{len(recs)+1}", "color": RGBColor(0xAB, 0x47, 0xBC),
        "title": "Celebrate & Share Positive Staff Feedback",
        "body": "Share named staff mentions from positive reviews with department heads. Nominate standout caregivers for the Care Champions program."
    })

    rec_colors_bg = RGBColor(0x0A, 0x20, 0x40)
    rec_border = RGBColor(0x1A, 0x3F, 0x6F)
    cols = 2 if len(recs) > 2 else 1
    w = 4.5 if cols == 2 else 9.0

    for i, r in enumerate(recs[:6]):
        col = i % cols
        row = i // cols
        x = 0.5 + col * (w + 0.2)
        y = 1.05 + row * 1.5

        add_rect(slide, x, y, w, 1.35, rec_colors_bg, rec_border, 0.5)
        add_textbox(slide, r["num"], x + 0.12, y + 0.1, 0.55, 0.55,
                    font_size=20, bold=True, color=r["color"])
        add_textbox(slide, r["title"], x + 0.72, y + 0.1, w - 0.88, 0.38,
                    font_size=10, bold=True, color=WHITE)
        add_textbox(slide, r["body"], x + 0.72, y + 0.5, w - 0.88, 0.76,
                    font_size=8.5, color=RGBColor(0xA0, 0xC4, 0xD8))


# ── MAIN ──────────────────────────────────────────────────────────────────────
def build_presentation(metrics: dict, output_path: str):
    meta = metrics["meta"]
    month_label = meta["month_label"]

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    print("🖼️  Building slides...")

    # Original 7-slide structure
    build_title_slide(prs, meta)
    build_overview_slide(prs, metrics, meta)
    build_regional_overview_slide(prs, metrics, meta)
    build_owner_slide(prs, metrics, meta)

    # Alerts — network then per region
    build_alerts_slide(prs, metrics, meta)
    for region in ["South", "North", "West"]:
        build_alerts_slide(prs, metrics, meta, region=region)

    # Recommendations — network then per region
    build_recommendations_slide(prs, metrics, meta)
    for region in ["South", "North", "West"]:
        build_recommendations_slide(prs, metrics, meta, region=region)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ Report saved to: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Build WMCHealth CTL PowerPoint report")
    parser.add_argument("--data", default="scripts/data.json", help="Path to metrics JSON from fetch_data.py")
    parser.add_argument("--output", default=None, help="Output PPTX path (default: docs/reports/WMCHealth_CTL_YYYY-MM.pptx)")
    args = parser.parse_args()

    with open(args.data, "r") as f:
        metrics = json.load(f)

    meta = metrics["meta"]
    month_str = meta["from_date"][:7]  # YYYY-MM

    output_path = args.output or f"docs/reports/WMCHealth_CTL_{month_str}.pptx"
    build_presentation(metrics, output_path)


if __name__ == "__main__":
    main()
